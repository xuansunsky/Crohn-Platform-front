# -*- coding: utf-8 -*-
"""
Crohn Paradise 融资路演 PPT v2 —— 投行级密度
设计：深色/暖白交替，Georgia 衬线做英文大标题与数字，雅黑做中文正文
统一页眉(kicker+标题+页码)与页脚(机密声明)栏，网格对齐
数据为演示用虚构数据
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 设计令牌 ----------
INK    = RGBColor(0x0C, 0x10, 0x18)
INK2   = RGBColor(0x15, 0x1B, 0x27)
CARDD  = RGBColor(0x1B, 0x22, 0x30)
LINE_D = RGBColor(0x2A, 0x33, 0x44)
PAPER  = RGBColor(0xF6, 0xF4, 0xEF)
CARDL  = RGBColor(0xFF, 0xFF, 0xFF)
LINE_L = RGBColor(0xE2, 0xDE, 0xD4)
INKT   = RGBColor(0x10, 0x16, 0x22)   # 暖白页上的深字
SUBL   = RGBColor(0x5B, 0x62, 0x70)   # 暖白页副文
CLOUD  = RGBColor(0x97, 0xA1, 0xB2)   # 深页副文
FAINT  = RGBColor(0x5C, 0x66, 0x78)   # 深页弱文
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x3B, 0x6E, 0xF5)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
EMER   = RGBColor(0x0F, 0xB8, 0x81)
AMBER  = RGBColor(0xE2, 0x8C, 0x09)
ROSE   = RGBColor(0xF4, 0x3F, 0x5E)

CN = "Microsoft YaHei"
SERIF = "Georgia"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
MARGIN = Inches(0.85)
CONTENT_W = SW - 2 * MARGIN


def slide(): return prs.slides.add_slide(BLANK)

def bg(s, c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def rect(s, x, y, w, h, color=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp

def text(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=1.0, sa=0):
    """lines: [[(t,size,color,bold,font), ...], ...]"""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if lines and not isinstance(lines[0], list):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        if sa: p.space_after = Pt(sa)
        for run in ln:
            t, sz, col, bold = run[0], run[1], run[2], run[3]
            fnt = run[4] if len(run) > 4 else CN
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = fnt
    return tb

def header(s, kicker, kcolor, title, pageno, dark=True):
    """统一页眉：kicker 细线标 + 大标题 + 右上页码 + 标题下细分隔线"""
    tcol = WHITE if dark else INKT
    sub = CLOUD if dark else SUBL
    rect(s, MARGIN, Inches(0.62), Inches(0.32), Inches(0.09), kcolor)
    text(s, MARGIN + Inches(0.45), Inches(0.5), Inches(8), Inches(0.4),
         [[(kicker, 11.5, kcolor, True)]])
    text(s, MARGIN, Inches(0.95), Inches(10.5), Inches(0.85),
         [[(title, 27, tcol, True)]])
    text(s, SW - MARGIN - Inches(1.6), Inches(0.55), Inches(1.6), Inches(0.4),
         [[(pageno, 11, sub, False, SERIF)]], align=PP_ALIGN.RIGHT)
    rect(s, MARGIN, Inches(1.78), CONTENT_W, Pt(1),
         LINE_D if dark else LINE_L)

def footer(s, dark=True):
    sub = FAINT if dark else RGBColor(0xA8, 0xA2, 0x96)
    rect(s, MARGIN, SH - Inches(0.62), CONTENT_W, Pt(0.75),
         LINE_D if dark else LINE_L)
    text(s, MARGIN, SH - Inches(0.55), Inches(8), Inches(0.35),
         [[("CROHN PARADISE", 9, sub, True, SERIF),
           ("   ·   种子轮路演   ·   严格保密", 9, sub, False)]])
    text(s, SW - MARGIN - Inches(3), SH - Inches(0.55), Inches(3), Inches(0.35),
         [[("Confidential — 2026", 9, sub, False, SERIF)]], align=PP_ALIGN.RIGHT)

def card(s, x, y, w, h, dark=True, accent=None):
    fill = CARDD if dark else CARDL
    ln = LINE_D if dark else LINE_L
    c = rect(s, x, y, w, h, fill, line=ln, lw=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.045
    if accent:
        rect(s, x, y + Inches(0.22), Inches(0.06), h - Inches(0.44), accent)
    return c


# ════════════════ 01 封面 ════════════════
s = slide(); bg(s, INK)
rect(s, 0, 0, SW, Inches(0.14), BLUE)
rect(s, 0, SH - Inches(0.14), SW, Inches(0.14), EMER)
# 顶栏
text(s, MARGIN, Inches(0.62), Inches(8), Inches(0.4),
     [[("CROHN PARADISE", 13, WHITE, True, SERIF),
       ("   /   患者共建的 IBD 互助社区", 12, CLOUD, False)]])
text(s, SW - MARGIN - Inches(3), Inches(0.64), Inches(3), Inches(0.4),
     [[("SEED ROUND · 2026", 11, CLOUD, False, SERIF)]], align=PP_ALIGN.RIGHT)
rect(s, MARGIN, Inches(1.15), CONTENT_W, Pt(0.75), LINE_D)
# 主标题
text(s, MARGIN, Inches(2.25), Inches(11.6), Inches(2.2),
     [[("把克罗恩病友", 50, WHITE, True)],
      [("从孤岛，连成", 50, WHITE, True), ("大陆", 50, BLUE, True)]], spacing=1.05)
text(s, MARGIN, Inches(4.55), Inches(11), Inches(0.6),
     [[("Connecting IBD patients — from isolated islands into one continent.", 15, CLOUD, False, SERIF)]])
# 底部三栏元信息
meta = [("赛道", "数字医疗 · 病种社区"), ("阶段", "产品上线在即 · 种子轮"), ("创始人", "林溪 · 患者 / 全栈")]
for i, (k, v) in enumerate(meta):
    x = MARGIN + i * Inches(3.95)
    rect(s, x, Inches(5.75), Inches(0.04), Inches(0.7), BLUE)
    text(s, x + Inches(0.2), Inches(5.72), Inches(3.6), Inches(0.35), [[(k, 11, FAINT, True)]])
    text(s, x + Inches(0.2), Inches(6.05), Inches(3.6), Inches(0.45), [[(v, 14.5, WHITE, True)]])
footer(s)


# ════════════════ 02 定位 ════════════════
s = slide(); bg(s, PAPER)
header(s, "POSITIONING / 一句话定位", BLUE, "我们是 IBD 病友的垂直生活方式平台", "02 — 12", dark=False)
text(s, MARGIN, Inches(2.2), Inches(11.8), Inches(1.4),
     [[("“", 30, BLUE, True, SERIF),
       ("大众点评的实测 ＋ 小红书的社区 ＋ 高德的医保导航", 29, INKT, True),
       ("”", 30, BLUE, True, SERIF)]], spacing=1.15)
text(s, MARGIN, Inches(3.7), Inches(11.6), Inches(1),
     [[("一个让克罗恩、溃疡性结肠炎患者", 17, SUBL, False),
       ("能安全说话、找到同类、交换保命信息", 17, INKT, True),
       ("的地方。", 17, SUBL, False)]], spacing=1.35)
# 三个差异点
diff = [("不是泛社区", "只服务 IBD 人群，高浓度、高信任"),
        ("不是问诊", "病友实测经验，不替代医疗、引导就医"),
        ("不是工具堆", "社区 × 数据 × 连接，三位一体")]
for i, (t, d) in enumerate(diff):
    x = MARGIN + i * Inches(3.95)
    card(s, x, Inches(5.0), Inches(3.7), Inches(1.55), dark=False, accent=BLUE)
    text(s, x + Inches(0.35), Inches(5.25), Inches(3.1), Inches(0.45), [[(t, 16, INKT, True)]])
    text(s, x + Inches(0.35), Inches(5.78), Inches(3.2), Inches(0.7), [[(d, 12.5, SUBL, False)]], spacing=1.2)
footer(s, dark=False)


# ════════════════ 03 痛点 ════════════════
s = slide(); bg(s, INK)
header(s, "THE PROBLEM / 痛点", ROSE, "在普通平台，病友是沉默的少数", "03 — 12")
intro = ("克罗恩病是终身、易复发的炎症性肠病。患者长期承受身体痛苦与社交隔离——"
         "而现有平台没有一个真正为他们设计。")
text(s, MARGIN, Inches(1.95), Inches(11.6), Inches(0.7), [[(intro, 14, CLOUD, False)]], spacing=1.3)
pains = [
    ("01", "说了没人懂", "在大众平台倾诉病情，常被无视或被健康人“多喝热水”式回应，造成二次心理伤害。", ROSE),
    ("02", "信息散又旧", "“某药本市能否报销”“哪家医院 IBD 门诊靠谱”——答案分散、过时、夹杂广告，决策成本极高。", AMBER),
    ("03", "深夜的孤独", "病痛发作、焦虑甚至绝望时，身边健康人难以共情，缺乏一个能即时连接同类的支撑网络。", VIOLET),
]
for i, (n, t, d, col) in enumerate(pains):
    x = MARGIN + i * Inches(3.95)
    card(s, x, Inches(2.85), Inches(3.7), Inches(3.35), accent=col)
    text(s, x + Inches(0.35), Inches(3.05), Inches(2), Inches(0.6), [[(n, 26, col, True, SERIF)]])
    text(s, x + Inches(0.35), Inches(3.75), Inches(3.1), Inches(0.5), [[(t, 18, WHITE, True)]])
    text(s, x + Inches(0.35), Inches(4.35), Inches(3.15), Inches(1.7), [[(d, 12.5, CLOUD, False)]], spacing=1.3)
footer(s)


# ════════════════ 04 解法 ════════════════
s = slide(); bg(s, PAPER)
header(s, "THE SOLUTION / 解法", EMER, "一个 App，四个高频“保命入口”", "04 — 12", dark=False)
feats = [
    ("01", "战术小队", "病友状态墙 · 密友圈动态 · 同城匹配 · 加密密聊与互助呼救", EMER),
    ("02", "医保政策地图", "城市级实测：门诊特殊病种 / 双通道药房 / 报销比例 / 凭证截图", BLUE),
    ("03", "药物图谱", "生物制剂实测安全度、起效反馈、双通道购药与价格情报", VIOLET),
    ("04", "经验金库", "带时间、地点、图片证据的病友手记，可信、可查、可沉淀", AMBER),
]
for i, (n, t, d, col) in enumerate(feats):
    x = MARGIN + (i % 2) * Inches(5.9)
    y = Inches(2.15) + (i // 2) * Inches(2.15)
    card(s, x, y, Inches(5.55), Inches(1.9), dark=False, accent=col)
    text(s, x + Inches(0.4), y + Inches(0.28), Inches(1.2), Inches(0.6), [[(n, 22, col, True, SERIF)]])
    text(s, x + Inches(1.35), y + Inches(0.3), Inches(4), Inches(0.5), [[(t, 18, INKT, True)]])
    text(s, x + Inches(1.35), y + Inches(0.92), Inches(4.0), Inches(0.85), [[(d, 12, SUBL, False)]], spacing=1.25)
footer(s, dark=False)


# ════════════════ 05 为什么是我们 ════════════════
s = slide(); bg(s, INK)
header(s, "WHY US / 不公平的优势", AMBER, "创始人本人，就是克罗恩患者", "05 — 12")
text(s, MARGIN, Inches(1.95), Inches(11.6), Inches(1.0),
     [[("这不是“寻找赛道”，而是“为自己和同类挖一口井”。病友能一眼认出你是自己人——", 15, CLOUD, False)],
      [("这种与生俱来的信任，是健康人创业者用任何预算都买不到的冷启动燃料。", 15, CLOUD, False)]], spacing=1.35)
moats = [
    ("患者身份", "真实共情转化为最低成本获客", "Patient-led trust"),
    ("全栈自研", "产品 / 后端 / 数据一人贯通，迭代极快", "Full-stack velocity"),
    ("数据壁垒", "城市级医保实测库越用越厚，难以复制", "Compounding data moat"),
]
for i, (t, d, en) in enumerate(moats):
    x = MARGIN + i * Inches(3.95)
    card(s, x, Inches(3.05), Inches(3.7), Inches(3.1), accent=AMBER)
    text(s, x + Inches(0.35), Inches(3.3), Inches(3.1), Inches(0.5), [[(t, 19, WHITE, True)]])
    text(s, x + Inches(0.35), Inches(3.95), Inches(3.15), Inches(1.4), [[(d, 13.5, CLOUD, False)]], spacing=1.3)
    rect(s, x + Inches(0.35), Inches(5.55), Inches(2.9), Pt(0.75), LINE_D)
    text(s, x + Inches(0.35), Inches(5.7), Inches(3.1), Inches(0.4), [[(en, 11, FAINT, False, SERIF)]])
footer(s)


# ════════════════ 06 市场 ════════════════
s = slide(); bg(s, PAPER)
header(s, "MARKET / 市场规模", BLUE, "小而精，但浓度极高的人群", "06 — 12", dark=False)
mk = [("≈ 10M", "中国 IBD 患者群体", "且呈持续上升趋势", BLUE),
      ("15–25%", "近十年患病率年增速", "增长最快的慢病之一", VIOLET),
      ("终身", "慢病管理周期", "= 高频、高复购需求", EMER)]
for i, (n, d, sub, col) in enumerate(mk):
    x = MARGIN + i * Inches(3.95)
    card(s, x, Inches(2.15), Inches(3.7), Inches(2.5), dark=False, accent=col)
    text(s, x + Inches(0.35), Inches(2.45), Inches(3.2), Inches(0.9), [[(n, 40, col, True, SERIF)]])
    text(s, x + Inches(0.35), Inches(3.55), Inches(3.2), Inches(0.45), [[(d, 14, INKT, True)]])
    text(s, x + Inches(0.35), Inches(4.05), Inches(3.2), Inches(0.45), [[(sub, 12, SUBL, False)]])
card(s, MARGIN, Inches(5.0), CONTENT_W, Inches(1.5), dark=False, accent=AMBER)
text(s, MARGIN + Inches(0.4), Inches(5.25), Inches(11), Inches(1.1),
     [[("战略判断：", 15, INKT, True),
       ("我们不与小红书比“大”。3,000 名真实活跃的 IBD 病友，其商业价值远高于 30 万泛用户——", 14, SUBL, False)],
      [("因为他们高需求、高信任、可精准触达，是医疗产业链上游真正稀缺的资源。", 14, SUBL, False)]], spacing=1.35)
footer(s, dark=False)


# ════════════════ 07 商业模式 ════════════════
s = slide(); bg(s, INK)
header(s, "BUSINESS MODEL / 商业模式", EMER, "社区是入口，数据与连接是收入", "07 — 12")
biz = [
    ("临床试验患者招募", "药企 / CRO 最缺合规患者线索，单条高价值", "短期", EMER),
    ("生物制剂 · DTP 导流", "长期高价复购药，厂商急需触达真实患者", "短期", BLUE),
    ("患者管理服务", "慢病随访依从性管理，药企愿为效果付费", "中期", VIOLET),
    ("增值会员 · 营养器械", "病友营养品、付费咨询、专科医生连接", "中期", AMBER),
]
for i, (t, d, term, col) in enumerate(biz):
    x = MARGIN + (i % 2) * Inches(5.9)
    y = Inches(2.15) + (i // 2) * Inches(2.05)
    card(s, x, y, Inches(5.55), Inches(1.8), accent=col)
    text(s, x + Inches(0.4), y + Inches(0.28), Inches(3.6), Inches(0.5), [[(t, 17, WHITE, True)]])
    rt = rect(s, x + Inches(4.35), y + Inches(0.3), Inches(0.85), Inches(0.36), None, line=col, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rt.adjustments[0] = 0.5
    tf = rt.text_frame; pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = term; rr.font.size = Pt(10.5); rr.font.bold = True; rr.font.color.rgb = col; rr.font.name = CN
    text(s, x + Inches(0.4), y + Inches(0.92), Inches(4.8), Inches(0.8), [[(d, 12.5, CLOUD, False)]], spacing=1.25)
footer(s)


# ════════════════ 08 早期数据 ════════════════
s = slide(); bg(s, PAPER)
header(s, "TRACTION / 早期数据", BLUE, "种子用户已经在“说话”", "08 — 12", dark=False)
tr = [("320+", "种子病友", BLUE), ("68%", "次周留存", EMER),
      ("1,400+", "实测经验条数", VIOLET), ("23", "覆盖城市医保数据", AMBER)]
for i, (n, d, col) in enumerate(tr):
    x = MARGIN + i * Inches(2.95)
    card(s, x, Inches(2.15), Inches(2.7), Inches(2.2), dark=False, accent=col)
    text(s, x + Inches(0.3), Inches(2.45), Inches(2.3), Inches(0.9), [[(n, 34, col, True, SERIF)]])
    text(s, x + Inches(0.3), Inches(3.5), Inches(2.4), Inches(0.6), [[(d, 12.5, SUBL, False)]], spacing=1.15)
card(s, MARGIN, Inches(4.7), CONTENT_W, Inches(1.75), dark=False, accent=EMER)
text(s, MARGIN + Inches(0.4), Inches(4.95), Inches(11), Inches(1.3),
     [[("关键信号 — ", 15, INKT, True),
       ("社区内已出现病友自发回答、互相“送暖”与经验沉淀，", 14, SUBL, False)],
      [("内容飞轮开始自转。这是病种社区从“工具”跨越到“网络”的第一道分水岭，", 14, SUBL, False)],
      [("也是早期最难、最值钱的信号。", 14, SUBL, False)]], spacing=1.35)
footer(s, dark=False)


# ════════════════ 09 增长路线 ════════════════
s = slide(); bg(s, INK)
header(s, "GO-TO-MARKET / 增长打法", AMBER, "先扎根，再蔓延", "09 — 12")
steps = [
    ("PHASE 1", "扎根", "创始人以病友身份，从病友群与医生处手动邀请前 30 名真实用户，并亲自填充种子内容，确保新人“60 秒内获得有用感”。", EMER),
    ("PHASE 2", "口碑", "城市级医保实测数据形成口碑资产，病友自发邀请同城战友，单城密度滚动提升。", BLUE),
    ("PHASE 3", "飞轮", "内容与数据飞轮成型，联合 IBD 公益组织与药企，覆盖全国主要城市，跑通商业闭环。", VIOLET),
]
y0 = Inches(2.2)
for i, (ph, t, d, col) in enumerate(steps):
    y = y0 + i * Inches(1.4)
    dot = rect(s, MARGIN + Inches(0.05), y + Inches(0.05), Inches(0.8), Inches(0.8), col, shape=MSO_SHAPE.OVAL)
    text(s, MARGIN + Inches(0.05), y + Inches(0.07), Inches(0.8), Inches(0.76),
         [[(str(i+1), 22, WHITE, True, SERIF)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        rect(s, MARGIN + Inches(0.42), y + Inches(0.88), Pt(1.2), Inches(0.5), LINE_D)
    text(s, MARGIN + Inches(1.2), y, Inches(2), Inches(0.35), [[(ph, 10.5, col, True, SERIF)]])
    text(s, MARGIN + Inches(1.2), y + Inches(0.32), Inches(2.1), Inches(0.55), [[(t, 22, WHITE, True)]])
    text(s, MARGIN + Inches(3.55), y + Inches(0.05), Inches(8.0), Inches(1.1),
         [[(d, 14, CLOUD, False)]], spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)
footer(s)


# ════════════════ 10 团队 ════════════════
s = slide(); bg(s, PAPER)
header(s, "TEAM / 团队", BLUE, "最小、最懂、最快的核心团队", "10 — 12", dark=False)
card(s, MARGIN, Inches(2.15), CONTENT_W, Inches(2.5), dark=False, accent=BLUE)
av = rect(s, MARGIN + Inches(0.45), Inches(2.55), Inches(1.7), Inches(1.7), BLUE, shape=MSO_SHAPE.OVAL)
text(s, MARGIN + Inches(0.45), Inches(2.55), Inches(1.7), Inches(1.7),
     [[("林", 44, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, MARGIN + Inches(2.55), Inches(2.6), Inches(8.5), Inches(0.55), [[("林溪", 24, INKT, True), ("   创始人 / 全栈工程师", 15, SUBL, False)]])
text(s, MARGIN + Inches(2.55), Inches(3.3), Inches(8.8), Inches(1.2),
     [[("克罗恩病友本人，5 年全栈研发经验，独立完成本产品前后端架构与数据体系。", 14, SUBL, False)],
      [("既是产品的第一个用户，也是它的建造者——理解痛点的深度，决定了产品的温度。", 14, SUBL, False)]], spacing=1.35)
card(s, MARGIN, Inches(4.95), CONTENT_W, Inches(1.55), dark=False, accent=AMBER)
text(s, MARGIN + Inches(0.4), Inches(5.2), Inches(11), Inches(0.45), [[("正在招募的关键合伙人", 14, INKT, True)]])
text(s, MARGIN + Inches(0.4), Inches(5.68), Inches(11), Inches(0.7),
     [[("社区运营合伙人", 13.5, BLUE, True), ("（病友社群增长）", 13, SUBL, False),
       ("      医疗 BD", 13.5, EMER, True), ("（药企 / CRO / 公益组织资源对接）", 13, SUBL, False)]], spacing=1.2)
footer(s, dark=False)


# ════════════════ 11 融资 ════════════════
s = slide(); bg(s, INK)
header(s, "THE ASK / 融资", EMER, "种子轮 ¥300 万，出让 10%", "11 — 12")
use = [("50%", "社区运营与种子病友扩张", EMER),
       ("30%", "医保 / 药物数据采集与核验", BLUE),
       ("20%", "合规建设与医疗 BD 拓展", AMBER)]
for i, (n, d, col) in enumerate(use):
    x = MARGIN + i * Inches(3.95)
    card(s, x, Inches(2.15), Inches(3.7), Inches(2.3), accent=col)
    text(s, x + Inches(0.35), Inches(2.45), Inches(3.2), Inches(0.9), [[(n, 36, col, True, SERIF)]])
    text(s, x + Inches(0.35), Inches(3.55), Inches(3.2), Inches(0.75), [[(d, 13.5, CLOUD, False)]], spacing=1.25)
# 里程碑条
card(s, MARGIN, Inches(4.75), CONTENT_W, Inches(1.7), accent=VIOLET)
text(s, MARGIN + Inches(0.4), Inches(5.0), Inches(11), Inches(0.45), [[("18 个月里程碑", 15, WHITE, True)]])
ms = [("50 城", "医保数据覆盖"), ("5 万", "真实活跃病友"), ("1 条", "药企付费闭环跑通")]
for i, (n, d) in enumerate(ms):
    x = MARGIN + Inches(0.4) + i * Inches(3.85)
    text(s, x, Inches(5.55), Inches(3.6), Inches(0.5), [[(n, 22, VIOLET, True, SERIF), ("  " + d, 13, CLOUD, False)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s)


# ════════════════ 12 收尾 ════════════════
s = slide(); bg(s, INK)
rect(s, 0, 0, SW, Inches(0.14), EMER)
rect(s, 0, SH - Inches(0.14), SW, Inches(0.14), BLUE)
text(s, MARGIN, Inches(0.7), Inches(8), Inches(0.4),
     [[("CLOSING", 11.5, EMER, True, SERIF)]])
text(s, MARGIN, Inches(2.35), Inches(11.8), Inches(2.0),
     [[("健康人有一万个平台，", 38, WHITE, True)],
      [("病友只需要", 38, WHITE, True), ("一个家。", 38, EMER, True)]], spacing=1.12)
text(s, MARGIN, Inches(4.55), Inches(11.5), Inches(0.6),
     [[("Crohn Paradise — connecting patients into one continent.", 15, CLOUD, False, SERIF)]])
rect(s, MARGIN, Inches(5.6), CONTENT_W, Pt(0.75), LINE_D)
contact = [("邮箱", "hi@crohnparadise.com"), ("微信", "crohn_paradise"), ("创始人", "林溪")]
for i, (k, v) in enumerate(contact):
    x = MARGIN + i * Inches(3.95)
    text(s, x, Inches(5.85), Inches(3.6), Inches(0.35), [[(k, 11, FAINT, True)]])
    text(s, x, Inches(6.2), Inches(3.6), Inches(0.45), [[(v, 15, WHITE, True)]])
text(s, SW - MARGIN - Inches(1.6), Inches(0.72), Inches(1.6), Inches(0.4),
     [[("12 — 12", 11, CLOUD, False, SERIF)]], align=PP_ALIGN.RIGHT)


out = r"f:\Dev\crohn-planform-front\CrohnParadise_Pitch_v2.pptx"
prs.save(out)
print("SAVED:", out, "| slides:", len(prs.slides._sldIdLst))
