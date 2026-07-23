# -*- coding: utf-8 -*-
"""克罗恩乐园 · 精品展示 PPT — 全中文、暖白+墨绿、无多余英文"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 色板：暖白纸感 + 森林绿 + 琥珀点缀
PAPER   = RGBColor(0xF7, 0xF2, 0xEA)
CREAM   = RGBColor(0xFF, 0xF9, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x19, 0x23, 0x1F)
MUTED   = RGBColor(0x69, 0x75, 0x6E)
SOFT    = RGBColor(0x8A, 0x76, 0x67)
MOSS    = RGBColor(0x50, 0x6F, 0x5B)
MOSS_L  = RGBColor(0xDD, 0xE8, 0xD7)
RIVER   = RGBColor(0x7B, 0xA6, 0xA6)
GOLD    = RGBColor(0xD8, 0xA7, 0x4E)
ROSE    = RGBColor(0xD9, 0x8A, 0x7A)
DARK    = RGBColor(0x10, 0x20, 0x1A)
DARK2   = RGBColor(0x1B, 0x36, 0x2E)
LINE_L  = RGBColor(0xE5, 0xDE, 0xD1)
LINE_D  = RGBColor(0x31, 0x56, 0x4B)
CLOUD   = RGBColor(0xB9, 0xCB, 0xBF)
PALE    = RGBColor(0xD7, 0xE6, 0xD7)
CN = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
M = Inches(0.78)
CW = SW - 2 * M


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c


def box(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, r=0.045):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = r
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, sa=0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if lines and not isinstance(lines[0], list):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if sa:
            p.space_after = Pt(sa)
        for run in ln:
            t, sz, col, bold = run[0], run[1], run[2], run[3]
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.name = CN
    return tb


def kicker(s, text, dark=False):
    col = MOSS_L if dark else MOSS
    txt(s, M, Inches(0.48), Inches(4), Inches(0.32), [[(text, 10.5, col, True)]])
    box(s, M, Inches(0.86), Inches(0.95), Pt(2.5), fill=col)


def pagenum(s, n, dark=False):
    col = CLOUD if dark else SOFT
    txt(s, SW - M - Inches(0.7), Inches(6.88), Inches(0.7), Inches(0.28),
        [[(f"{n:02d}", 9, col, True)]], align=PP_ALIGN.RIGHT)


def title(s, lines, dark=False, size=34, y=Inches(1.12)):
    col = WHITE if dark else INK
    if isinstance(lines, str):
        lines = [lines]
    txt(s, M, y, Inches(10.8), Inches(1.6),
        [[(ln, size, col, True)] for ln in lines], spacing=1.08)


def card(s, x, y, w, h, dark=False, accent=None):
    fill = DARK2 if dark else WHITE
    ln = LINE_D if dark else LINE_L
    box(s, x, y, w, h, fill=fill, line=ln, lw=0.75)
    if accent:
        box(s, x, y + Inches(0.18), Inches(0.055), h - Inches(0.36), fill=accent, line=accent)


def pill(s, text, x, y, w, dark=False):
    fill = RGBColor(0x24, 0x4D, 0x43) if dark else CREAM
    ln = MOSS_L if dark else LINE_L
    box(s, x, y, w, Inches(0.38), fill=fill, line=ln, lw=0.75, r=0.18)
    txt(s, x + Inches(0.14), y + Inches(0.08), w - Inches(0.28), Inches(0.22),
        [[(text, 9.5, PALE if dark else SOFT, True)]], align=PP_ALIGN.CENTER)


def backdrop_light(s):
    bg(s, PAPER)
    box(s, Inches(9.2), Inches(-0.8), Inches(4.5), Inches(4.5), fill=MOSS_L, line=MOSS_L, shape=MSO_SHAPE.OVAL)
    box(s, Inches(-1.5), Inches(5.2), Inches(4.2), Inches(2.8), fill=RGBColor(0xE1, 0xEF, 0xF0), line=RGBColor(0xE1, 0xEF, 0xF0), shape=MSO_SHAPE.OVAL)


def backdrop_dark(s):
    bg(s, DARK)
    box(s, Inches(8.5), Inches(-1.2), Inches(5.5), Inches(5.5), fill=DARK2, line=LINE_D, lw=1.2, shape=MSO_SHAPE.OVAL)
    box(s, Inches(-1.0), Inches(4.5), Inches(5.0), Inches(3.2), fill=RGBColor(0x20, 0x3A, 0x36), line=LINE_D, lw=0.8, shape=MSO_SHAPE.OVAL)


# ── 01 封面 ──
s = slide()
backdrop_dark(s)
pill(s, "7月1日 · 小范围内测", M, Inches(0.62), Inches(2.15), dark=True)
title(s, ["克罗恩乐园", "病友互助工具箱"], dark=True, size=44, y=Inches(1.55))
txt(s, M, Inches(3.05), Inches(6.2), Inches(0.9),
    [[("记经验 · 查医保 · 看食评 · 找同伴", 17, PALE, True)]])
txt(s, M, Inches(3.95), Inches(5.8), Inches(0.55),
    [[("不是热闹广场，是一块低噪音、能慢慢活下去的小地方。", 14, CLOUD, False)]], spacing=1.35)
for i, (k, v) in enumerate([("定位", "病友工具箱"), ("阶段", "内测打磨"), ("态度", "真实反馈优先")]):
    x = M + i * Inches(3.6)
    txt(s, x, Inches(5.55), Inches(3.2), Inches(0.3), [[(k, 10, CLOUD, True)]])
    txt(s, x, Inches(5.88), Inches(3.2), Inches(0.4), [[(v, 15, WHITE, True)]])
box(s, Inches(8.8), Inches(1.35), Inches(3.4), Inches(5.6), fill=RGBColor(0x15, 0x2A, 0x24), line=LINE_D, lw=1.0)
txt(s, Inches(9.15), Inches(1.75), Inches(2.8), Inches(0.3), [[("今日你能做的事", 12, MOSS_L, True)]])
for i, t in enumerate(["记录身体反馈", "查同城医保", "看病友经验", "找懂你的人"]):
    txt(s, Inches(9.15), Inches(2.35 + i * 0.72), Inches(2.8), Inches(0.35),
        [[("·  " + t, 12.5, RGBColor(0xEE, 0xF6, 0xE9), False)]])


# ── 02 定位 ──
s = slide()
backdrop_light(s)
kicker(s, "01 / 一句话定位")
pagenum(s, 2)
title(s, ["给克罗恩病友用的互助工具"], size=32)
txt(s, M, Inches(2.15), Inches(11.5), Inches(1.1),
    [[("“", 28, GOLD, True), ("实测经验", 26, INK, True),
      (" + ", 26, MUTED, False), ("同城连接", 26, INK, True),
      (" + ", 26, MUTED, False), ("医保情报", 26, INK, True), ("”", 28, GOLD, True)]], spacing=1.2)
items = [
    ("不是问诊", "病友经验分享，不替代医生", MOSS),
    ("不是广场", "低噪音、有边界、能沉淀", RIVER),
    ("不是工具堆", "围绕真实生活线串联", GOLD),
]
for i, (t, d, c) in enumerate(items):
    x = M + i * Inches(3.85)
    card(s, x, Inches(3.55), Inches(3.55), Inches(2.35), accent=c)
    txt(s, x + Inches(0.35), Inches(3.85), Inches(3), Inches(0.45), [[(t, 17, INK, True)]])
    txt(s, x + Inches(0.35), Inches(4.45), Inches(3.1), Inches(0.9), [[(d, 12.5, MUTED, False)]], spacing=1.3)


# ── 03 痛点 ──
s = slide()
backdrop_dark(s)
kicker(s, "02 / 为什么需要", dark=True)
pagenum(s, 3, dark=True)
title(s, ["经验散在聊天里，", "真正需要时却找不到"], dark=True, size=32)
pains = [
    ("01", "说了没人懂", "在大众平台倾诉，常被无视或「多喝热水」，造成二次伤害。", ROSE),
    ("02", "信息散又旧", "医保、医院、用药反馈分散在群聊和朋友圈，难以检索。", GOLD),
    ("03", "深夜的孤独", "病痛发作时，身边健康人难共情，缺一个能连接同类的支撑网。", RIVER),
]
for i, (n, t, d, c) in enumerate(pains):
    x = M + i * Inches(3.85)
    card(s, x, Inches(2.85), Inches(3.55), Inches(3.45), dark=True, accent=c)
    txt(s, x + Inches(0.35), Inches(3.05), Inches(1.2), Inches(0.55), [[(n, 24, c, True)]])
    txt(s, x + Inches(0.35), Inches(3.75), Inches(3), Inches(0.45), [[(t, 18, WHITE, True)]])
    txt(s, x + Inches(0.35), Inches(4.35), Inches(3.05), Inches(1.6), [[(d, 12, CLOUD, False)]], spacing=1.3)


# ── 04 五模块 ──
s = slide()
backdrop_light(s)
kicker(s, "03 / 五个核心模块")
pagenum(s, 4)
title(s, ["五个界面，围绕病友一天的真实问题"], size=31)
mods = [
    ("战术小队", "聊天、密友、同城、呼救", MOSS, "队"),
    ("社区情报", "食物实测、身体反馈", ROSE, "评"),
    ("医保地图", "城市政策、医院入口", RIVER, "图"),
    ("经验金库", "手记、标签、地点、主题", GOLD, "库"),
    ("我的主页", "资料、签名、个人动态", SOFT, "我"),
]
xs = [M + i * Inches(2.38) for i in range(5)]
for i, (name, desc, col, badge) in enumerate(mods):
    y = Inches(2.05) if i % 2 == 0 else Inches(2.55)
    dark_card = i == 0
    card(s, xs[i], y, Inches(2.15), Inches(3.35), dark=dark_card, accent=col)
    av = box(s, xs[i] + Inches(0.32), y + Inches(0.32), Inches(0.62), Inches(0.62),
             fill=col if not dark_card else RGBColor(0x31, 0x56, 0x4B),
             line=col, shape=MSO_SHAPE.OVAL)
    txt(s, xs[i] + Inches(0.32), y + Inches(0.42), Inches(0.62), Inches(0.42),
        [[(badge, 11, WHITE if dark_card else INK, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tc = WHITE if dark_card else INK
    dc = CLOUD if dark_card else MUTED
    txt(s, xs[i] + Inches(0.28), y + Inches(1.25), Inches(1.6), Inches(0.35), [[(name, 15, tc, True)]])
    txt(s, xs[i] + Inches(0.28), y + Inches(1.85), Inches(1.55), Inches(0.75), [[(desc, 9.5, dc, False)]], spacing=1.2)


# ── 05 经验金库 ──
s = slide()
backdrop_light(s)
kicker(s, "04 / 经验金库亮点")
pagenum(s, 5)
title(s, ["小红书式布局，成熟标签体系"], size=31)
feats = [
    ("内容标签", "求助 · 经验分享 · 知识分享 · 交友", MOSS),
    ("话题标签", "用药、饮食、医保、医院、复查…", RIVER),
    ("地点标记", "认证城市 + 可选具体地点", GOLD),
    ("阅读体验", "列表有风格，详情可一键「简洁」", ROSE),
]
for i, (t, d, c) in enumerate(feats):
    x = M + (i % 2) * Inches(5.85)
    y = Inches(2.35) + (i // 2) * Inches(1.85)
    card(s, x, y, Inches(5.55), Inches(1.55), accent=c)
    txt(s, x + Inches(0.38), y + Inches(0.32), Inches(4.8), Inches(0.38), [[(t, 16, INK, True)]])
    txt(s, x + Inches(0.38), y + Inches(0.82), Inches(4.9), Inches(0.5), [[(d, 12.5, MUTED, False)]])
card(s, Inches(8.65), Inches(2.15), Inches(3.55), Inches(4.35), dark=True, accent=GOLD)
txt(s, Inches(9.0), Inches(2.55), Inches(2.9), Inches(0.3), [[("设计理念", 12, MOSS_L, True)]])
for i, t in enumerate(["列表：卡片有性格", "详情：默认奶白专业", "沉浸：尊重发帖人", "简洁：阅读者可关闭"]):
    txt(s, Inches(9.0), Inches(3.15 + i * 0.62), Inches(2.9), Inches(0.35),
        [[("·  " + t, 11.5, RGBColor(0xEE, 0xF6, 0xE9), False)]])


# ── 06 医保地图 ──
s = slide()
backdrop_dark(s)
kicker(s, "05 / 医保政策地图", dark=True)
pagenum(s, 6, dark=True)
title(s, ["城市实测，办理前先查"], dark=True, size=32)
metrics = [("23+", "覆盖城市"), ("门特", "特殊病种"), ("双通道", "购药通道"), ("12393", "核实热线")]
for i, (n, d) in enumerate(metrics):
    x = M + i * Inches(2.95)
    card(s, x, Inches(2.35), Inches(2.65), Inches(2.15), dark=True, accent=RIVER if i % 2 else MOSS)
    txt(s, x + Inches(0.3), Inches(2.65), Inches(2.1), Inches(0.75), [[(n, 30, MOSS_L, True)]])
    txt(s, x + Inches(0.3), Inches(3.55), Inches(2.2), Inches(0.4), [[(d, 12, CLOUD, False)]])
txt(s, M, Inches(4.85), Inches(11.5), Inches(1.2),
    [[("病友上传报销凭证与城市实测，让后来者少跑冤枉路。", 15, CLOUD, False)],
     [("数据来源标注清楚，办理前仍建议致电当地 12393 核实。", 14, CLOUD, False)]], spacing=1.35)


# ── 07 设计原则 ──
s = slide()
backdrop_light(s)
kicker(s, "06 / 设计原则")
pagenum(s, 7)
title(s, ["克制、温暖、有空气感"], size=32)
principles = [
    ("少一点装饰", "页面先服务判断，不抢注意力"),
    ("多一点呼吸", "卡片间距让信息能被看懂"),
    ("不要程序员味", "去掉无意义英文和调试感"),
    ("隐私可感知", "能不能被找到，用户自己决定"),
]
for i, (t, d) in enumerate(principles):
    x = M if i < 2 else M + Inches(6.05)
    y = Inches(2.45) + (i % 2) * Inches(1.45)
    card(s, x, y, Inches(5.55), Inches(1.15), accent=[MOSS, RIVER, ROSE, GOLD][i])
    txt(s, x + Inches(0.38), y + Inches(0.28), Inches(4.8), Inches(0.32), [[(t, 15, INK, True)]])
    txt(s, x + Inches(0.38), y + Inches(0.68), Inches(4.9), Inches(0.28), [[(d, 11, MUTED, False)]])


# ── 08 内测计划 ──
s = slide()
backdrop_light(s)
kicker(s, "07 / 7月1日内测")
pagenum(s, 8)
title(s, ["先找 30 个真实病友，", "把小闭环打磨到愿意留下"], size=31)
phases = [
    ("第 1 批", "5–10 人", "验证登录、聊天、食评、医保入口", MOSS),
    ("第 2 批", "20–30 人", "收集哪里难用、哪里多余、哪里有用", RIVER),
    ("持续修", "听反馈", "只加高频刚需，不被大而全拖死", GOLD),
]
for i, (ph, num, desc, c) in enumerate(phases):
    x = M + i * Inches(3.85)
    card(s, x, Inches(2.95), Inches(3.55), Inches(2.85), dark=(i == 1), accent=c)
    tc = WHITE if i == 1 else INK
    dc = CLOUD if i == 1 else MUTED
    txt(s, x + Inches(0.35), Inches(3.25), Inches(1.2), Inches(0.28), [[(ph, 10.5, dc, True)]])
    txt(s, x + Inches(0.35), Inches(3.62), Inches(2.8), Inches(0.45), [[(num, 22, tc, True)]])
    txt(s, x + Inches(0.35), Inches(4.25), Inches(2.95), Inches(0.9), [[(desc, 11.5, dc, False)]], spacing=1.25)


# ── 09 邀请话术 ──
s = slide()
backdrop_dark(s)
kicker(s, "08 / 内测邀请", dark=True)
pagenum(s, 9, dark=True)
title(s, ["我也是克罗恩，", "自己写了个小东西"], dark=True, size=36)
card(s, M, Inches(3.05), Inches(7.35), Inches(2.45), dark=True, accent=MOSS)
txt(s, M + Inches(0.45), Inches(3.45), Inches(6.5), Inches(1.6),
    [[("专门给咱们这种人记经验、查医保、看食物反馈、找同城病友。", 16, RGBColor(0xEE, 0xF6, 0xE9), True)],
     [("7月1日开很小的内测版。难用你直接骂我，有用也请告诉我。", 15, CLOUD, False)]], spacing=1.35)
for i, t in enumerate(["能不能登录", "能不能看懂", "有没有真用", "哪里最烦"]):
    pill(s, t, Inches(9.05), Inches(2.55 + i * 0.58), Inches(2.35), dark=True)


# ── 10 边界 ──
s = slide()
backdrop_light(s)
kicker(s, "09 / 我们坚持什么")
pagenum(s, 10)
title(s, ["诚实发布，比假装成熟更重要"], size=31)
left = [("不做医疗诊断", MOSS), ("不替代医生建议", RIVER), ("不把隐私当流量", ROSE), ("不做无限刷人", GOLD)]
for i, (t, c) in enumerate(left):
    y = Inches(2.35) + i * Inches(0.72)
    txt(s, M, y, Inches(0.25), Inches(0.28), [[("—", 16, c, True)]])
    txt(s, M + Inches(0.35), y + Inches(0.02), Inches(4.5), Inches(0.28), [[(t, 14, INK, True)]])
card(s, Inches(6.85), Inches(2.15), Inches(5.35), Inches(3.55), accent=MOSS)
txt(s, Inches(7.2), Inches(2.55), Inches(4.6), Inches(0.35), [[("我们坚持", 14, INK, True)]])
for i, t in enumerate(["真实反馈比漂亮口号重要", "小功能闭环比大而空重要", "隐私边界比热闹重要", "能帮到一个人，就值得继续修"]):
    txt(s, Inches(7.2), Inches(3.15 + i * 0.58), Inches(4.6), Inches(0.35),
        [[("·  " + t, 12, MUTED, False)]])


# ── 11 体验路径 ──
s = slide()
backdrop_dark(s)
kicker(s, "10 / 首版体验路径", dark=True)
pagenum(s, 11, dark=True)
title(s, ["从「今天不舒服」", "到「我知道下一步怎么做」"], dark=True, size=31)
steps = [("记录", "身体状态、病例、经验"), ("查询", "医保、地图、可用信息"),
         ("判断", "看食物实测，不盲目模仿"), ("连接", "找到能理解你的人")]
for i, (st, desc) in enumerate(steps):
    x = M + Inches(0.1) + i * Inches(2.95)
    av = box(s, x, Inches(4.15), Inches(0.55), Inches(0.55), fill=MOSS_L, line=MOSS_L, shape=MSO_SHAPE.OVAL)
    txt(s, x, Inches(4.22), Inches(0.55), Inches(0.42),
        [[(str(i + 1), 12, DARK, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        box(s, x + Inches(0.62), Inches(4.38), Inches(2.05), Pt(1.5), fill=LINE_D, line=LINE_D)
    txt(s, x - Inches(0.15), Inches(4.95), Inches(1.0), Inches(0.32), [[(st, 15, WHITE, True)]], align=PP_ALIGN.CENTER)
    txt(s, x - Inches(0.55), Inches(5.38), Inches(1.8), Inches(0.45),
        [[(desc, 9.5, CLOUD, False)]], align=PP_ALIGN.CENTER)


# ── 12 收尾 ──
s = slide()
backdrop_dark(s)
box(s, 0, 0, SW, Inches(0.12), fill=MOSS, line=MOSS)
box(s, 0, SH - Inches(0.12), SW, Inches(0.12), fill=GOLD, line=GOLD)
txt(s, M, Inches(2.15), Inches(11.5), Inches(1.8),
    [[("先让它对一小部分人", 38, WHITE, True)],
     [("真的有用。", 38, MOSS_L, True)]], spacing=1.1)
txt(s, M, Inches(4.05), Inches(8.5), Inches(0.45),
    [[("如果能帮到一个人，就值得继续修。", 16, CLOUD, False)]])
box(s, M, Inches(5.15), Inches(2.4), Pt(2), fill=MOSS_L, line=MOSS_L)
txt(s, M, Inches(5.55), Inches(5), Inches(0.28), [[("克罗恩乐园 · 7月1日内测", 12, CLOUD, True)]])
for i, (k, v) in enumerate([("产品", "克罗恩乐园"), ("阶段", "内测招募"), ("联系", "欢迎真实病友反馈")]):
    x = M + i * Inches(3.6)
    txt(s, x, Inches(6.15), Inches(3.2), Inches(0.28), [[(k, 10, RGBColor(0x88, 0xA0, 0x90), True)]])
    txt(s, x, Inches(6.48), Inches(3.2), Inches(0.35), [[(v, 14, WHITE, True)]])


out = r"F:\Dev\crohn-planform-front\CrohnParadise_Showcase.pptx"
prs.save(out)
print("SAVED:", out, "| slides:", len(prs.slides))
